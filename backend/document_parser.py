"""
Multi-format document parser: PPTX, XLSX, PDF, DOCX, TXT, JSON, XML, CSV, MD.
Extends utils.process_resume() with presentation, spreadsheet, and structured data support.
"""

import csv
import json
import os
import xml.etree.ElementTree as ET


def parse_pptx(file_path):
    """Extract text from PowerPoint slides including tables."""
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"--- Slide {i} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        lines.append(" | ".join(cells))
        return "\n".join(lines)
    except Exception as e:
        return f"[PPTX parse error: {e}]"


def parse_xlsx(file_path):
    """Extract text from all Excel sheets."""
    try:
        from openpyxl import load_workbook

        wb = load_workbook(file_path, read_only=True, data_only=True)
        lines = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append(f"--- Sheet: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                non_empty = [c for c in cells if c]
                if non_empty:
                    lines.append(" | ".join(cells))
        wb.close()
        return "\n".join(lines)
    except Exception as e:
        return f"[XLSX parse error: {e}]"


def parse_json(file_path):
    """Extract readable text from JSON files (config, data, exports)."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            data = json.load(f)
        return _json_to_text(data)
    except Exception as e:
        return f"[JSON parse error: {e}]"


def _json_to_text(obj, prefix=""):
    """Recursively flatten JSON into readable key-value text."""
    lines = []
    if isinstance(obj, dict):
        for key, val in obj.items():
            path = f"{prefix}.{key}" if prefix else key
            if isinstance(val, (dict, list)):
                lines.append(_json_to_text(val, path))
            else:
                lines.append(f"{path}: {val}")
    elif isinstance(obj, list):
        for i, item in enumerate(obj):
            path = f"{prefix}[{i}]"
            if isinstance(item, (dict, list)):
                lines.append(_json_to_text(item, path))
            else:
                lines.append(f"{path}: {item}")
    else:
        lines.append(f"{prefix}: {obj}" if prefix else str(obj))
    return "\n".join(lines)


def parse_xml(file_path):
    """Extract text content from XML documents."""
    try:
        tree = ET.parse(file_path)
        root = tree.getroot()
        lines = []
        _xml_walk(root, lines)
        return "\n".join(lines)
    except Exception as e:
        return f"[XML parse error: {e}]"


def _xml_walk(element, lines, depth=0):
    """Recursively extract text from XML elements."""
    tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag
    text = (element.text or "").strip()
    tail = (element.tail or "").strip()

    if text:
        lines.append(f"{tag}: {text}")
    elif element.attrib:
        attrs = ", ".join(f"{k}={v}" for k, v in element.attrib.items())
        lines.append(f"{tag} ({attrs})")

    for child in element:
        _xml_walk(child, lines, depth + 1)

    if tail:
        lines.append(tail)


def parse_csv(file_path):
    """Extract text from CSV files as pipe-delimited rows."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            lines = []
            for row in reader:
                non_empty = [c for c in row if c.strip()]
                if non_empty:
                    lines.append(" | ".join(row))
        return "\n".join(lines)
    except Exception as e:
        return f"[CSV parse error: {e}]"


def parse_md(file_path):
    """Read markdown files as-is — already AI-readable."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()
    except Exception as e:
        return f"[Markdown parse error: {e}]"


def parse_any(file_path):
    """Route by extension, fallback to utils.process_resume()."""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pptx":
        return parse_pptx(file_path)
    elif ext in (".xlsx", ".xls"):
        return parse_xlsx(file_path)
    elif ext == ".json":
        return parse_json(file_path)
    elif ext == ".xml":
        return parse_xml(file_path)
    elif ext == ".csv":
        return parse_csv(file_path)
    elif ext == ".md":
        return parse_md(file_path)
    else:
        # Use existing parsers for PDF, DOCX, DOC, TXT
        try:
            from utils import process_resume

            result = process_resume(file_path)
            if isinstance(result, dict):
                return result.get("text", str(result))
            return str(result)
        except Exception as e:
            # Last resort: try reading as text
            try:
                with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                    return f.read()
            except Exception:
                return f"[Parse error: {e}]"
